"""Build PPTX slide describing the Analyst's workflow.

Mirrors the visual style of the Post-Payment AI Powered Detection slide:
- Left column: title block, callout box, checkmark bullets
- Right column: 5 numbered step tiles with circle indicators and timeline
- Palette: navy / teal / PayGuard pink
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ── Palette (matches the reference slide) ─────────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x5F)
TEAL        = RGBColor(0x5A, 0x9B, 0x96)
PINK        = RGBColor(0xFE, 0x01, 0x7D)   # PayGuard accent
BLACK       = RGBColor(0x11, 0x11, 0x11)
DARK_GRAY   = RGBColor(0x3D, 0x3D, 0x3D)
MED_GRAY    = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY  = RGBColor(0xE5, 0xE7, 0xEB)
CALLOUT_BG  = RGBColor(0xE6, 0xED, 0xF5)
CALLOUT_BD  = RGBColor(0xCB, 0xD7, 0xE8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)


def set_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb: RGBColor, width_pt: float = 0.75) -> None:
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)


def no_line(shape) -> None:
    line = shape.line
    line.fill.background()


def add_text(
    shape, text: str, *,
    size: int = 12, bold: bool = False, color: RGBColor = BLACK,
    align=PP_ALIGN.LEFT, font: str = "Calibri",
    anchor=MSO_ANCHOR.TOP, margin: float = 0.08,
) -> None:
    tf = shape.text_frame
    tf.margin_left   = Inches(margin)
    tf.margin_right  = Inches(margin)
    tf.margin_top    = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # Replace existing paragraph
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_multi_text(
    shape, parts: list, *,
    align=PP_ALIGN.LEFT, font: str = "Calibri",
    anchor=MSO_ANCHOR.TOP, margin: float = 0.08,
) -> None:
    """parts: list of dicts with keys text, size, bold, color, and optional 'new_para'."""
    tf = shape.text_frame
    tf.margin_left   = Inches(margin)
    tf.margin_right  = Inches(margin)
    tf.margin_top    = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    for part in parts:
        if part.get("new_para"):
            p = tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(part.get("space_before", 0))
        run = p.add_run()
        run.text = part["text"]
        run.font.name = font
        run.font.size = Pt(part.get("size", 12))
        run.font.bold = part.get("bold", False)
        run.font.color.rgb = part.get("color", BLACK)


# ── Build slide ───────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# White background (default), but draw a soft pink accent in top-right corner
accent_blob = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(11.6), Inches(-0.6), Inches(2.5), Inches(2.5)
)
set_fill(accent_blob, RGBColor(0xFD, 0xE9, 0xF2))
no_line(accent_blob)

# Top-left pink bar (small)
top_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.25), Inches(0.6), Inches(0.06)
)
set_fill(top_bar, PINK)
no_line(top_bar)


# ── Left column: title + description ──────────────────────────────────────

title1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(4.6), Inches(0.85))
add_text(title1, "Analyst Workflow", size=36, bold=True, color=BLACK, margin=0.0)

title2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.6), Inches(0.75))
add_text(title2, "Human-in-the-Loop Recovery", size=26, bold=True, color=PINK, margin=0.0)

desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.95), Inches(4.6), Inches(1.3))
add_text(
    desc,
    "Analysts review AI-prioritized cases, adjudicate findings, "
    "issue provider notices, and reconcile recoveries — every action "
    "auditable and reversible.",
    size=14, color=DARK_GRAY, margin=0.0,
)


# ── Left column: Glass-Box-style callout ──────────────────────────────────

callout = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(3.55), Inches(4.6), Inches(2.05),
)
set_fill(callout, CALLOUT_BG)
set_line(callout, CALLOUT_BD, 0.75)
callout.adjustments[0] = 0.04

callout_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(4.2), Inches(0.5))
add_text(
    callout_title, "Glass Box — Full Accountability",
    size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER, margin=0.0,
)

callout_body = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(4.0), Inches(1.4))
add_multi_text(
    callout_body,
    [
        {"text": "Decisions traced to evidence",        "size": 13, "color": DARK_GRAY},
        {"text": "Per-finding accept / reject / adjust", "size": 13, "color": DARK_GRAY, "new_para": True, "space_before": 4},
        {"text": "Notes, @mentions, contact log",       "size": 13, "color": DARK_GRAY, "new_para": True, "space_before": 4},
        {"text": "Reversible & version-controlled",     "size": 13, "color": DARK_GRAY, "new_para": True, "space_before": 4},
    ],
    margin=0.0,
)


# ── Left column: checkmark bullets ────────────────────────────────────────

bullets = [
    "Faster Throughput",
    "Higher Recovery Yield",
    "Defensible Decisions",
]
by = 5.95
for label in bullets:
    # Green check circle
    chk = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(by), Inches(0.22), Inches(0.22))
    set_fill(chk, GREEN)
    no_line(chk)
    add_text(chk, "✓", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    # Label
    lbl = slide.shapes.add_textbox(Inches(0.8), Inches(by - 0.04), Inches(4.0), Inches(0.32))
    add_text(lbl, label, size=14, bold=True, color=PINK, margin=0.0)
    by += 0.4


# ── Right column: vertical timeline + step tiles ──────────────────────────

steps = [
    {
        "num": "01", "color": NAVY,
        "title": "Triage & Pickup",
        "body": "Reviews the prioritized worklist, scopes to my-cases or all-cases, "
                "and picks up the top-ranked items. AI has already classified and ranked them.",
    },
    {
        "num": "02", "color": NAVY,
        "title": "Evidence Review",
        "body": "Inspects firing detectors, claim lines, member and provider context, prior "
                "audit log, and provider risk attribution to build a complete picture.",
    },
    {
        "num": "03", "color": TEAL,
        "title": "Disposition & Adjustment",
        "body": "Accepts, rejects, or adjusts each finding individually; overrides the "
                "at-risk amount when warranted; documents reasoning in @mentioned notes.",
    },
    {
        "num": "04", "color": PINK,
        "title": "Notice & Escalation",
        "body": "Generates the provider notice from templates, sends it via mail or fax, "
                "and escalates to a supervisor when the case exceeds the $2K threshold.",
    },
    {
        "num": "05", "color": PINK,
        "title": "Reconciliation & Closure",
        "body": "Logs provider contact, matches incoming 835 reversals to expected recoupment, "
                "and records the final disposition — recovered, written off, or overturned.",
    },
]

right_x = Inches(5.6)
tile_x  = Inches(6.5)
tile_w  = Inches(6.4)
tile_h  = Inches(1.05)
circle_d = Inches(0.62)
top_y = 0.45
row_gap = 1.25

# Timeline vertical line — drawn first so circles sit on top
line_x = right_x + circle_d / 2
line_top = Inches(top_y + 0.62)
line_bot = Inches(top_y + (len(steps) - 1) * row_gap)
timeline = slide.shapes.add_connector(1, line_x, line_top, line_x, line_bot)
timeline.line.color.rgb = LIGHT_GRAY
timeline.line.width = Pt(1.25)

for i, step in enumerate(steps):
    y = Inches(top_y + i * row_gap)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x, y, circle_d, circle_d)
    set_fill(circle, step["color"])
    no_line(circle)
    add_text(
        circle, step["num"],
        size=12, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.0,
    )

    # Tile card
    tile = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        tile_x, y - Inches(0.18), tile_w, tile_h,
    )
    set_fill(tile, WHITE)
    set_line(tile, LIGHT_GRAY, 0.75)
    tile.adjustments[0] = 0.08

    # Colored left edge accent on the tile
    edge = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        tile_x, y - Inches(0.18), Inches(0.05), tile_h,
    )
    set_fill(edge, step["color"])
    no_line(edge)

    # Tile text (title + body)
    txt_x = tile_x + Inches(0.25)
    txt_w = tile_w - Inches(0.45)
    txt = slide.shapes.add_textbox(txt_x, y - Inches(0.08), txt_w, tile_h - Inches(0.2))
    add_multi_text(
        txt,
        [
            {"text": step["title"], "size": 15, "bold": True, "color": BLACK},
            {"text": step["body"], "size": 11, "color": DARK_GRAY, "new_para": True, "space_before": 2},
        ],
        margin=0.0,
    )


# ── Save ──────────────────────────────────────────────────────────────────

out = "/Users/issamzeinoun/claude/overcoding/opa/analyst_workflow.pptx"
prs.save(out)
print(f"Wrote {out}")
