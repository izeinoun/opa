"""Build a 3-slide deck describing PayGuard's technical capabilities for a
senior payment-integrity executive audience. Mid-level between functional
and technical. Shares the visual language of build_workflow_deck.py.

   1. Multi-Layer Detection Engine
   2. Explainable Risk Scoring & Prioritization
   3. AI-Augmented Analyst Workflow with Audit-Grade Traceability
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Palette (matches build_workflow_deck.py) ──────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x5F)
TEAL        = RGBColor(0x5A, 0x9B, 0x96)
PINK        = RGBColor(0xFE, 0x01, 0x7D)
BLACK       = RGBColor(0x11, 0x11, 0x11)
DARK_GRAY   = RGBColor(0x3D, 0x3D, 0x3D)
LIGHT_GRAY  = RGBColor(0xE5, 0xE7, 0xEB)
CALLOUT_BG  = RGBColor(0xE6, 0xED, 0xF5)
CALLOUT_BD  = RGBColor(0xCB, 0xD7, 0xE8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def add_text(shape, text, *, size=12, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, font="Calibri",
             anchor=MSO_ANCHOR.TOP, margin=0.08):
    tf = shape.text_frame
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", Inches(margin))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_multi_text(shape, parts, *, align=PP_ALIGN.LEFT, font="Calibri",
                   anchor=MSO_ANCHOR.TOP, margin=0.08):
    tf = shape.text_frame
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", Inches(margin))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    for part in parts:
        if part.get("new_para"):
            p = tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(part.get("space_before", 0))
        run = p.add_run()
        run.text = part["text"]
        run.font.name = font
        run.font.size = Pt(part.get("size", 12))
        run.font.bold = part.get("bold", False)
        run.font.color.rgb = part.get("color", BLACK)


# ── Slide builder ─────────────────────────────────────────────────────────

def build_slide(prs, *, title_main, title_sub, description,
                callout_title, callout_bullets,
                check_bullets, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top-right pink blob
    blob = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(11.6), Inches(-0.6), Inches(2.5), Inches(2.5)
    )
    set_fill(blob, RGBColor(0xFD, 0xE9, 0xF2))
    no_line(blob)

    # Top-left pink bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.25), Inches(0.6), Inches(0.06)
    )
    set_fill(bar, PINK)
    no_line(bar)

    # ── Left column ───────────────────────────────────────────────────────
    t1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(4.6), Inches(0.85))
    add_text(t1, title_main, size=34, bold=True, color=BLACK, margin=0.0)

    t2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.6), Inches(0.75))
    add_text(t2, title_sub, size=24, bold=True, color=PINK, margin=0.0)

    desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.95), Inches(4.6), Inches(1.5))
    add_text(desc, description, size=13, color=DARK_GRAY, margin=0.0)

    # Callout box
    co = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.55), Inches(4.6), Inches(2.15),
    )
    set_fill(co, CALLOUT_BG)
    set_line(co, CALLOUT_BD, 0.75)
    co.adjustments[0] = 0.04

    co_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(4.2), Inches(0.5))
    add_text(co_title, callout_title, size=15, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, margin=0.0)

    co_body = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(4.0), Inches(1.5))
    parts = []
    for i, b in enumerate(callout_bullets):
        parts.append({
            "text": b, "size": 12, "color": DARK_GRAY,
            "new_para": i > 0, "space_before": 4,
        })
    add_multi_text(co_body, parts, margin=0.0)

    # Check bullets
    by = 6.05
    for label in check_bullets:
        chk = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(by),
                                     Inches(0.22), Inches(0.22))
        set_fill(chk, GREEN)
        no_line(chk)
        add_text(chk, "✓", size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
        lbl = slide.shapes.add_textbox(Inches(0.8), Inches(by - 0.04),
                                       Inches(4.0), Inches(0.32))
        add_text(lbl, label, size=13, bold=True, color=PINK, margin=0.0)
        by += 0.38

    # ── Right column ──────────────────────────────────────────────────────
    right_x = Inches(5.6)
    tile_x  = Inches(6.5)
    tile_w  = Inches(6.4)
    tile_h  = Inches(1.15)
    circle_d = Inches(0.62)
    top_y = 0.45
    row_gap = 1.35

    # Timeline
    line_x = right_x + circle_d / 2
    line_top = Inches(top_y + 0.62)
    line_bot = Inches(top_y + (len(steps) - 1) * row_gap)
    timeline = slide.shapes.add_connector(1, line_x, line_top, line_x, line_bot)
    timeline.line.color.rgb = LIGHT_GRAY
    timeline.line.width = Pt(1.25)

    for i, step in enumerate(steps):
        y = Inches(top_y + i * row_gap)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x, y, circle_d, circle_d)
        set_fill(circle, step["color"])
        no_line(circle)
        add_text(circle, step["num"], size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)

        tile = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            tile_x, y - Inches(0.18), tile_w, tile_h,
        )
        set_fill(tile, WHITE)
        set_line(tile, LIGHT_GRAY, 0.75)
        tile.adjustments[0] = 0.08

        edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            tile_x, y - Inches(0.18), Inches(0.05), tile_h,
        )
        set_fill(edge, step["color"])
        no_line(edge)

        txt_x = tile_x + Inches(0.25)
        txt_w = tile_w - Inches(0.45)
        txt = slide.shapes.add_textbox(txt_x, y - Inches(0.08), txt_w,
                                       tile_h - Inches(0.2))
        add_multi_text(txt, [
            {"text": step["title"], "size": 14, "bold": True, "color": BLACK},
            {"text": step["body"],  "size": 10.5, "color": DARK_GRAY,
             "new_para": True, "space_before": 2},
        ], margin=0.0)


# ── Compose the deck ──────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Slide 1: Multi-Layer Detection Engine ─────────────────────────────────
build_slide(
    prs,
    title_main="Detection Engine",
    title_sub="Six Concerns. One Pass.",
    description=(
        "Every paid claim is screened against six purpose-built detectors "
        "spanning the most material overpayment vectors in commercial and "
        "government books — deterministic rules where the policy is bright-line, "
        "machine learning where billing behavior is the signal."
    ),
    callout_title="Why Six, Not One Model",
    callout_bullets=[
        "Each detector maps to a distinct regulatory or contractual lever",
        "Deterministic rules give defensible, auditable findings",
        "ML supplements, it never overrides, the rule-based logic",
        "New detectors plug into the same orchestrator without rework",
    ],
    check_bullets=[
        "Rule + ML hybrid",
        "Explainable per-finding evidence",
        "Configurable thresholds per LOB",
    ],
    steps=[
        {"num": "01", "color": NAVY, "title": "Duplicate Billing (DET-01)",
         "body": "Catches same member + CPT + service-date resubmissions across "
                 "providers and TINs — the single largest source of recoverable dollars."},
        {"num": "02", "color": NAVY, "title": "Retro Eligibility (DET-02)",
         "body": "Reconciles paid services against enrollment history to identify "
                 "claims paid for members not eligible on the date of service."},
        {"num": "04", "color": TEAL, "title": "Fee Schedule Variance (DET-04)",
         "body": "Flags allowed-vs-paid drift above contract tolerance — surfaces "
                 "configuration errors in the claim adjudication system before they compound."},
        {"num": "06", "color": TEAL, "title": "NCCI / MUE Violations (DET-06)",
         "body": "Applies CMS edit tables for mutually exclusive procedures and "
                 "medically unlikely units — clinical compliance, encoded."},
        {"num": "08", "color": PINK, "title": "Excluded Provider (DET-08)",
         "body": "Cross-references OIG/SAM exclusion lists against rendering and "
                 "billing NPIs at claim time — eliminates manual list reconciliation."},
        {"num": "09", "color": PINK, "title": "Coding & Dx Errors (DET-09)",
         "body": "Detects invalid ICD→CPT pairings and unbundling patterns — the "
                 "category most often missed by single-vendor edit engines."},
    ],
)


# ── Slide 2: Explainable Risk Scoring & Prioritization ────────────────────
build_slide(
    prs,
    title_main="Risk Scoring",
    title_sub="Explainable. Prioritized. Defensible.",
    description=(
        "Every flagged claim receives a composite likelihood score and a "
        "priority band that puts the highest-yield, time-sensitive recoveries "
        "at the top of the analyst queue — with the math fully exposed so "
        "findings hold up in provider appeals."
    ),
    callout_title="Glass-Box, Not Black-Box",
    callout_bullets=[
        "Likelihood = CPT risk × 0.30 + Provider tier × 0.25 + Dx/CPT × 0.20 + Complexity × 0.15 + ML variance × 0.10",
        "Priority blends amount-at-risk, likelihood, and deadline urgency",
        "SHAP attribution explains every provider-level risk score in plain English",
        "Thresholds, weights, and bands are operator-tunable — no model retrain required",
    ],
    check_bullets=[
        "Per-claim, per-provider transparency",
        "Tunable weights without redeploy",
        "Auditor-grade documentation",
    ],
    steps=[
        {"num": "01", "color": NAVY, "title": "Five-Factor Likelihood Model",
         "body": "Combines CPT risk, provider risk tier, Dx/CPT coherence, claim "
                 "complexity, and an AutoML billing-variance signal into a single 0–1 score."},
        {"num": "02", "color": NAVY, "title": "AutoML Billing Variance",
         "body": "Penguin FDEAutoML trained on seven behavioral features — unit "
                 "intensity, modifier patterns, peer deviation — outputs a per-provider score."},
        {"num": "03", "color": TEAL, "title": "Priority Score & Banding",
         "body": "Priority = (amount × 0.60 + likelihood × 0.35 + urgency × 0.05) × 100. "
                 "HIGH ≥ 75, MEDIUM 50–74, LOW < 50; ≤ 5 days to deadline forces HIGH."},
        {"num": "04", "color": TEAL, "title": "SHAP-Driven Explainability",
         "body": "Per-provider SHAP attribution rendered as plain-English narrative — "
                 "\"modifier usage is significantly elevated vs. peer specialty.\""},
        {"num": "05", "color": PINK, "title": "Operator-Controlled Tuning",
         "body": "Weights, thresholds, urgency windows, and detector rules are stored "
                 "in config — adjust for LOB or contract terms without touching code."},
    ],
)


# ── Slide 3: AI-Augmented Analyst Workflow ────────────────────────────────
build_slide(
    prs,
    title_main="Analyst Workflow",
    title_sub="AI-Augmented. Audit-Ready.",
    description=(
        "Claude Sonnet 4.6 — accessed through AWS Bedrock — drafts the case "
        "narrative, summarizes evidence, and authors provider letters. Every "
        "prompt, model response, rule fire, and human action is captured in "
        "an immutable timeline behind each case."
    ),
    callout_title="Defensible AI at Production Scale",
    callout_bullets=[
        "AWS Bedrock keeps PHI inside the customer's cloud perimeter",
        "Langfuse traces every LLM call — model, latency, tokens, cost, prompt",
        "Per-case audit log captures actor, action, status transitions, reasoning",
        "Analyst remains the decision-maker; AI accelerates, never overrides",
    ],
    check_bullets=[
        "HIPAA-aware deployment",
        "Full chain-of-custody",
        "Closed-loop recoupment",
    ],
    steps=[
        {"num": "01", "color": NAVY, "title": "Prioritized Worklist",
         "body": "Analysts see only high-priority, evidence-backed cases — sorted by "
                 "dollar yield and deadline pressure, with the detection rationale inline."},
        {"num": "02", "color": NAVY, "title": "AI-Drafted Case Narrative",
         "body": "Claude synthesizes claim facts, detector findings, member and provider "
                 "context into a reviewable summary — analyst edits, never starts from blank."},
        {"num": "03", "color": TEAL, "title": "Letter Generation & Disposition",
         "body": "Provider/UIC letters are auto-drafted from approved templates with "
                 "case-specific evidence; disposition codes drive the recoupment workflow."},
        {"num": "04", "color": TEAL, "title": "Immutable Audit Timeline",
         "body": "Every status change, override, comment, and AI generation is logged "
                 "with actor and timestamp — discovery-ready without forensic reconstruction."},
        {"num": "05", "color": PINK, "title": "Recoupment & Reconciliation",
         "body": "Closes the loop: tracks recovered dollars against findings, feeds "
                 "outcomes back into model retraining and provider risk tiers."},
    ],
)


# ── Output ────────────────────────────────────────────────────────────────
out = "/Users/issamzeinoun/claude/overcoding/opa/capabilities_deck.pptx"
prs.save(out)
print(f"Wrote {out}")
