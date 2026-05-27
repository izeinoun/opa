"""Build a 3-slide deck describing the PayGuard workflows:
   1. AI Detection pipeline (recreates the source slide)
   2. Analyst workflow
   3. Supervisor workflow

All slides share the same visual language (navy/teal/pink, rounded tiles,
numbered timeline, glass-box callout, green-check bullets).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Palette ───────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x5F)
TEAL        = RGBColor(0x5A, 0x9B, 0x96)
PINK        = RGBColor(0xFE, 0x01, 0x7D)
BLACK       = RGBColor(0x11, 0x11, 0x11)
DARK_GRAY   = RGBColor(0x3D, 0x3D, 0x3D)
LIGHT_GRAY  = RGBColor(0xE5, 0xE7, 0xEB)
CALLOUT_BG  = RGBColor(0xE6, 0xED, 0xF5)
CALLOUT_BD  = RGBColor(0xCB, 0xD7, 0xE8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def add_text(shape, text, *, size=12, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, font="Calibri",
             anchor=MSO_ANCHOR.TOP, margin=0.08):
    tf = shape.text_frame
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", Inches(margin))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_multi_text(shape, parts, *, align=PP_ALIGN.LEFT, font="Calibri",
                   anchor=MSO_ANCHOR.TOP, margin=0.08):
    tf = shape.text_frame
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", Inches(margin))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    for part in parts:
        if part.get("new_para"):
            p = tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(part.get("space_before", 0))
        run = p.add_run()
        run.text = part["text"]
        run.font.name = font
        run.font.size = Pt(part.get("size", 12))
        run.font.bold = part.get("bold", False)
        run.font.color.rgb = part.get("color", BLACK)


# ── Slide builder ─────────────────────────────────────────────────────────

def build_slide(prs, *, title_main, title_sub, description,
                callout_title, callout_bullets,
                check_bullets, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top-right pink blob
    blob = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(11.6), Inches(-0.6), Inches(2.5), Inches(2.5)
    )
    set_fill(blob, RGBColor(0xFD, 0xE9, 0xF2))
    no_line(blob)

    # Top-left pink bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.25), Inches(0.6), Inches(0.06)
    )
    set_fill(bar, PINK)
    no_line(bar)

    # ── Left column ───────────────────────────────────────────────────────
    t1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(4.6), Inches(0.85))
    add_text(t1, title_main, size=36, bold=True, color=BLACK, margin=0.0)

    t2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.6), Inches(0.75))
    add_text(t2, title_sub, size=26, bold=True, color=PINK, margin=0.0)

    desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.95), Inches(4.6), Inches(1.3))
    add_text(desc, description, size=14, color=DARK_GRAY, margin=0.0)

    # Callout box
    co = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.55), Inches(4.6), Inches(2.05),
    )
    set_fill(co, CALLOUT_BG)
    set_line(co, CALLOUT_BD, 0.75)
    co.adjustments[0] = 0.04

    co_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(4.2), Inches(0.5))
    add_text(co_title, callout_title, size=16, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, margin=0.0)

    co_body = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(4.0), Inches(1.4))
    parts = []
    for i, b in enumerate(callout_bullets):
        parts.append({
            "text": b, "size": 13, "color": DARK_GRAY,
            "new_para": i > 0, "space_before": 4,
        })
    add_multi_text(co_body, parts, margin=0.0)

    # Check bullets
    by = 5.95
    for label in check_bullets:
        chk = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(by),
                                     Inches(0.22), Inches(0.22))
        set_fill(chk, GREEN)
        no_line(chk)
        add_text(chk, "✓", size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
        lbl = slide.shapes.add_textbox(Inches(0.8), Inches(by - 0.04),
                                       Inches(4.0), Inches(0.32))
        add_text(lbl, label, size=14, bold=True, color=PINK, margin=0.0)
        by += 0.4

    # ── Right column ──────────────────────────────────────────────────────
    right_x = Inches(5.6)
    tile_x  = Inches(6.5)
    tile_w  = Inches(6.4)
    tile_h  = Inches(1.05)
    circle_d = Inches(0.62)
    top_y = 0.45
    row_gap = 1.25

    # Timeline
    line_x = right_x + circle_d / 2
    line_top = Inches(top_y + 0.62)
    line_bot = Inches(top_y + (len(steps) - 1) * row_gap)
    timeline = slide.shapes.add_connector(1, line_x, line_top, line_x, line_bot)
    timeline.line.color.rgb = LIGHT_GRAY
    timeline.line.width = Pt(1.25)

    for i, step in enumerate(steps):
        y = Inches(top_y + i * row_gap)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x, y, circle_d, circle_d)
        set_fill(circle, step["color"])
        no_line(circle)
        add_text(circle, step["num"], size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)

        tile = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            tile_x, y - Inches(0.18), tile_w, tile_h,
        )
        set_fill(tile, WHITE)
        set_line(tile, LIGHT_GRAY, 0.75)
        tile.adjustments[0] = 0.08

        edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            tile_x, y - Inches(0.18), Inches(0.05), tile_h,
        )
        set_fill(edge, step["color"])
        no_line(edge)

        txt_x = tile_x + Inches(0.25)
        txt_w = tile_w - Inches(0.45)
        txt = slide.shapes.add_textbox(txt_x, y - Inches(0.08), txt_w,
                                       tile_h - Inches(0.2))
        add_multi_text(txt, [
            {"text": step["title"], "size": 15, "bold": True, "color": BLACK},
            {"text": step["body"],  "size": 11, "color": DARK_GRAY,
             "new_para": True, "space_before": 2},
        ], margin=0.0)


# ── Compose the deck ──────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# Slide 1: AI Detection (recreated from the source image)
build_slide(
    prs,
    title_main="Post-Payment",
    title_sub="AI Powered Detection",
    description=(
        "Runs regularly with a look-back and rechecks 835s to detect "
        "incorrect payments — especially overpayments."
    ),
    callout_title="Glass Box — Full Transparency",
    callout_bullets=[
        "Applies known Rules",
        "Tracks leading indicators and reasoning",
        "Records intermediary and final outputs",
        "Logs auditable activity at each step",
    ],
    check_bullets=["AI Pipeline and Rules", "User Workflows", "Model & App Administration"],
    steps=[
        {"num": "01", "color": NAVY, "title": "Post-Payment Ingestion & Filtering",
         "body": "PayGuard ingests the 835s, runs a classification model, and selects "
                 "the most suspect payments for more specific detection rules."},
        {"num": "02", "color": NAVY, "title": "Detection & Scoring",
         "body": "PayGuard runs a mix of deterministic and LLM rules against each claim, "
                 "further classifying, prioritizing, and analyzing the claims prior to "
                 "Analyst reviews."},
        {"num": "03", "color": TEAL, "title": "Analyst Review & Recovery",
         "body": "A payment integrity analyst reviews the evidence-backed case, adjusts, "
                 "and actions on it: corrects, annotates, escalates. Sent to provider or UIC."},
        {"num": "04", "color": PINK, "title": "Reconciliation & Closure",
         "body": "PayGuard receives reversal transactions and matches the expected "
                 "recovery amount. Confirms and tracks financial closure, and records the "
                 "final disposition — recovered, written off, or overturned."},
        {"num": "05", "color": PINK, "title": "Autonomous Learning and Retraining",
         "body": "PayGuard regularly re-trains using new data from the adjudicated claim "
                 "and expert notes."},
    ],
)


# Slide 2: Analyst workflow
build_slide(
    prs,
    title_main="Analyst Workflow",
    title_sub="Human-in-the-Loop Recovery",
    description=(
        "Analysts review AI-prioritized cases, adjudicate findings, "
        "issue provider notices, and reconcile recoveries — every action "
        "auditable and reversible."
    ),
    callout_title="Glass Box — Full Accountability",
    callout_bullets=[
        "Decisions traced to evidence",
        "Per-finding accept / reject / adjust",
        "Notes, @mentions, contact log",
        "Reversible & version-controlled",
    ],
    check_bullets=["Faster Throughput", "Higher Recovery Yield", "Defensible Decisions"],
    steps=[
        {"num": "01", "color": NAVY, "title": "Triage & Pickup",
         "body": "Reviews the prioritized worklist, scopes to my-cases or all-cases, "
                 "and picks up the top-ranked items. AI has already classified and ranked them."},
        {"num": "02", "color": NAVY, "title": "Evidence Review",
         "body": "Inspects firing detectors, claim lines, member and provider context, "
                 "prior audit log, and provider risk attribution to build a complete picture."},
        {"num": "03", "color": TEAL, "title": "Disposition & Adjustment",
         "body": "Accepts, rejects, or adjusts each finding individually; overrides the "
                 "at-risk amount when warranted; documents reasoning in @mentioned notes."},
        {"num": "04", "color": PINK, "title": "Notice & Escalation",
         "body": "Generates the provider notice from templates, sends it via mail or fax, "
                 "and escalates to a supervisor when the case exceeds the $2K threshold."},
        {"num": "05", "color": PINK, "title": "Reconciliation & Closure",
         "body": "Logs provider contact, matches incoming 835 reversals to expected "
                 "recoupment, and records the final disposition — recovered, written off, "
                 "or overturned."},
    ],
)


# Slide 3: Supervisor workflow
build_slide(
    prs,
    title_main="Supervisor Workflow",
    title_sub="Oversight & Governance",
    description=(
        "Supervisors gate high-impact cases, triage escalations, balance "
        "team workload, and tune detection rules and prioritization weights."
    ),
    callout_title="Control Tower — Configurable Governance",
    callout_bullets=[
        "$2K approval gate on closures",
        "Escalation queue with audit trail",
        "Bulk assign across the analyst pool",
        "Tunable rules, weights, and thresholds",
    ],
    check_bullets=["Quality Control", "Capacity Planning", "Configurable Detection"],
    steps=[
        {"num": "01", "color": NAVY, "title": "Approval Queue ($2K Gate)",
         "body": "Reviews case closures flagged by the $2K threshold. Approves, rejects, "
                 "or returns to the analyst with a written decision and audit-logged rationale."},
        {"num": "02", "color": NAVY, "title": "Escalation Triage",
         "body": "Acts on soft-escalated cases — re-assigns, takes ownership, or resolves. "
                 "Each escalation carries its reason and notifies the originator on resolution."},
        {"num": "03", "color": TEAL, "title": "Bulk Workload Assignment",
         "body": "Selects multiple cases from the worklist and assigns them across the "
                 "analyst pool in a single action. Drives per-analyst throughput visibility."},
        {"num": "04", "color": PINK, "title": "Rule & Threshold Tuning",
         "body": "Adjusts prioritization weights, detector rule scores, supervisor gate "
                 "threshold, and letter templates — every change versioned and auditable."},
        {"num": "05", "color": PINK, "title": "Model & Team Performance",
         "body": "Retrains the billing-variance model on demand, monitors precision / recall / "
                 "AUC, reviews per-provider SHAP attributions, and tracks team performance."},
    ],
)


out = "/Users/issamzeinoun/claude/overcoding/opa/payguard_workflows.pptx"
prs.save(out)
print(f"Wrote {out}")
